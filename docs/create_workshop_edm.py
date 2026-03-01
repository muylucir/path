#!/usr/bin/env python3
"""P.A.T.H Workshop EDM Generator - A4+ vertical format with timeline"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x38, 0x8B, 0xF7)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_PINK = RGBColor(0xEC, 0x48, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)
CARD_BORDER = RGBColor(0x33, 0x41, 0x55)


def add_rrect(slide, left, top, width, height, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.05
    return s


def add_circle(slide, left, top, size, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def add_rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, left, top, width, height, text, size=12,
        color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return tb


def mtxt(slide, left, top, width, height, lines):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    for i, ln in enumerate(lines):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.text = ln.get("t", "")
        p.font.size = Pt(ln.get("s", 12))
        p.font.color.rgb = ln.get("c", WHITE)
        p.font.bold = ln.get("b", False)
        p.font.name = "Arial"
        p.alignment = ln.get("a", PP_ALIGN.LEFT)
        p.space_after = Pt(ln.get("sa", 2))
    return tb


def badge_text(shape, text, size=8, color=WHITE, bold=True):
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(size)
    shape.text_frame.paragraphs[0].font.color.rgb = color
    shape.text_frame.paragraphs[0].font.bold = bold
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


# ─── Dimensions ───
SW = Inches(7.5)
SH = Inches(14.5)  # A4+ to fit timeline
MG = Inches(0.45)
CW = SW - 2 * MG


def create_edm():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

    y = Inches(0)

    # ── Top accent bar ──
    add_rect(slide, 0, 0, SW, Inches(0.04), ACCENT_BLUE)

    # Deco circles
    add_circle(slide, Inches(5.6), Inches(0.2), Inches(0.9), RGBColor(0x1A, 0x3A, 0x5C))
    add_circle(slide, Inches(6.0), Inches(0.6), Inches(0.5), RGBColor(0x2A, 0x1F, 0x4E))

    # ══════════════════════════════════════
    # HERO
    # ══════════════════════════════════════
    y = Inches(0.35)

    b = add_rrect(slide, MG, y, Inches(2.2), Inches(0.28), RGBColor(0x1E, 0x3A, 0x5F), ACCENT_BLUE)
    badge_text(b, "  HANDS-ON WORKSHOP", 8, ACCENT_BLUE)

    y += Inches(0.45)

    mtxt(slide, MG, y, Inches(5.5), Inches(1.1), [
        {"t": "AI Agent,", "s": 28, "b": True, "sa": 0},
        {"t": "어디서부터 시작할까?", "s": 28, "b": True, "sa": 6},
        {"t": "P.A.T.H로 아이디어 검증부터 코드 구현까지 한번에", "s": 12, "c": ACCENT_CYAN},
    ])

    y += Inches(1.35)

    mtxt(slide, MG, y, CW, Inches(0.55), [
        {"t": "AI Agent를 도입하고 싶지만 어디서부터 시작해야 할지 모르겠다면,", "s": 10, "c": LIGHT_GRAY, "sa": 1},
        {"t": "이 워크샵에서 아이디어 검증 → 명세서 생성 → 코드 구현을 직접 경험하세요.", "s": 10, "c": WHITE, "b": True},
    ])

    y += Inches(0.65)
    add_rect(slide, MG, y, CW, Inches(0.015), CARD_BORDER)
    y += Inches(0.3)

    # ══════════════════════════════════════
    # P.A.T.H FRAMEWORK
    # ══════════════════════════════════════
    txt(slide, MG, y, CW, Inches(0.2), "P.A.T.H Framework", 8, ACCENT_BLUE, True)
    y += Inches(0.22)
    txt(slide, MG, y, CW, Inches(0.3), "4단계로 완성하는 AI Agent 설계", 17, WHITE, True)
    y += Inches(0.45)

    steps = [
        ("P", "Problem", "기본 정보 입력", ACCENT_BLUE),
        ("A", "Assessment", "준비도 점검", ACCENT_PURPLE),
        ("T", "Technical Review", "패턴 분석", ACCENT_CYAN),
        ("H", "Handoff", "명세서 생성", ACCENT_GREEN),
    ]

    cw = Inches(1.42)
    gap = Inches(0.14)

    for i, (letter, title, desc, color) in enumerate(steps):
        x = MG + i * (cw + gap)
        add_rrect(slide, x, y, cw, Inches(1.15), DARK_CARD, CARD_BORDER)

        c = add_circle(slide, x + Inches(0.41), y + Inches(0.12), Inches(0.45), color)
        badge_text(c, letter, 18, WHITE)

        txt(slide, x, y + Inches(0.65), cw, Inches(0.2), title, 9, color, True, PP_ALIGN.CENTER)
        txt(slide, x, y + Inches(0.85), cw, Inches(0.2), desc, 8, LIGHT_GRAY, False, PP_ALIGN.CENTER)

        if i < 3:
            txt(slide, x + cw + Inches(0.01), y + Inches(0.3), gap, Inches(0.25),
                ">", 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    y += Inches(1.4)

    # ══════════════════════════════════════
    # WORKSHOP FLOW (inline compact)
    # ══════════════════════════════════════
    txt(slide, MG, y, CW, Inches(0.2), "Workshop Flow", 8, ACCENT_PURPLE, True)
    y += Inches(0.22)
    txt(slide, MG, y, CW, Inches(0.3), "검증에서 구현까지, 한 자리에서", 17, WHITE, True)
    y += Inches(0.5)

    # Single row: Step A → Step B
    half = (CW - Inches(0.5)) / 2  # space for arrow in middle

    # Step A card
    add_rrect(slide, MG, y, half, Inches(0.9), DARK_CARD, ACCENT_BLUE)
    ba = add_rrect(slide, MG + Inches(0.15), y + Inches(0.12), Inches(0.7), Inches(0.2), ACCENT_BLUE)
    badge_text(ba, "Step A", 7, WHITE)
    txt(slide, MG + Inches(0.15), y + Inches(0.38), half - Inches(0.3), Inches(0.22),
        "PATH로 아이디어 검증", 11, WHITE, True)
    txt(slide, MG + Inches(0.15), y + Inches(0.6), half - Inches(0.3), Inches(0.22),
        "준비도 평가 + 패턴 분석 + 명세서", 8, ACCENT_BLUE)

    # Arrow
    arrow_x = MG + half
    txt(slide, arrow_x, y + Inches(0.25), Inches(0.5), Inches(0.4),
        ">>>", 14, LIGHT_GRAY, True, PP_ALIGN.CENTER)

    # Step B card
    bx = MG + half + Inches(0.5)
    add_rrect(slide, bx, y, half, Inches(0.9), DARK_CARD, ACCENT_GREEN)
    bb = add_rrect(slide, bx + Inches(0.15), y + Inches(0.12), Inches(0.7), Inches(0.2), ACCENT_GREEN)
    badge_text(bb, "Step B", 7, WHITE)
    txt(slide, bx + Inches(0.15), y + Inches(0.38), half - Inches(0.3), Inches(0.22),
        "AI-DLC로 코드 구현", 11, WHITE, True)
    txt(slide, bx + Inches(0.15), y + Inches(0.6), half - Inches(0.3), Inches(0.22),
        "동작 가능한 Agent 코드", 8, ACCENT_GREEN)

    y += Inches(1.15)

    # ══════════════════════════════════════
    # WHAT YOU'LL GET (2x2 grid)
    # ══════════════════════════════════════
    txt(slide, MG, y, CW, Inches(0.2), "What You'll Get", 8, ACCENT_GREEN, True)
    y += Inches(0.22)
    txt(slide, MG, y, CW, Inches(0.3), "워크샵에서 가져가는 것", 17, WHITE, True)
    y += Inches(0.45)

    gets = [
        ("01", "AI Agent 아이디어 검증 경험", "아이디어를 구조화하고 실현 가능성 평가", ACCENT_BLUE),
        ("02", "실제 동작하는 Agent 코드", "명세서 기반 AI-DLC가 생성한 코드", ACCENT_PURPLE),
        ("03", "본인 업무용 설계서", "자유 주제로 작성한 본인만의 명세서", ACCENT_GREEN),
        ("04", "재사용 가능한 프레임워크", "어떤 업무든 P.A.T.H로 분석하는 방법론", ACCENT_ORANGE),
    ]

    hw = (CW - Inches(0.12)) / 2

    for i, (num, title, desc, color) in enumerate(gets):
        col = i % 2
        row = i // 2
        x = MG + col * (hw + Inches(0.12))
        ry = y + row * Inches(0.82)

        add_rrect(slide, x, ry, hw, Inches(0.72), DARK_CARD, CARD_BORDER)

        nc = add_rrect(slide, x + Inches(0.12), ry + Inches(0.15), Inches(0.3), Inches(0.3), color)
        badge_text(nc, num, 10, WHITE)

        txt(slide, x + Inches(0.52), ry + Inches(0.1), hw - Inches(0.6), Inches(0.2),
            title, 10, WHITE, True)
        txt(slide, x + Inches(0.52), ry + Inches(0.35), hw - Inches(0.6), Inches(0.3),
            desc, 8, LIGHT_GRAY)

    y += Inches(1.8)

    # ══════════════════════════════════════
    # WHO SHOULD ATTEND
    # ══════════════════════════════════════
    txt(slide, MG, y, CW, Inches(0.2), "Who Should Attend", 8, ACCENT_CYAN, True)
    y += Inches(0.22)
    txt(slide, MG, y, CW, Inches(0.3), "이런 분들을 위한 워크샵", 17, WHITE, True)
    y += Inches(0.42)

    audiences = [
        "AI Agent 도입을 원하지만 어디서부터 시작할지 모르는 분",
        "PoC부터 시작해서 시간을 낭비한 경험이 있는 분",
        "명세서부터 코드 구현까지 End-to-End 경험을 원하는 분",
    ]

    for t in audiences:
        ck = add_rrect(slide, MG + Inches(0.1), y + Inches(0.02), Inches(0.2), Inches(0.2), ACCENT_CYAN)
        badge_text(ck, "v", 8, WHITE)
        txt(slide, MG + Inches(0.42), y, Inches(5.5), Inches(0.25), t, 10, WHITE)
        y += Inches(0.32)

    y += Inches(0.25)

    # ══════════════════════════════════════
    # TIMELINE
    # ══════════════════════════════════════
    txt(slide, MG, y, CW, Inches(0.2), "Timeline", 8, ACCENT_PINK, True)
    y += Inches(0.22)
    txt(slide, MG, y, CW, Inches(0.3), "워크샵 일정 (5시간)", 17, WHITE, True)
    y += Inches(0.5)

    timeline = [
        ("09:00", "09:20", "오프닝", "PATH 소개 & 워크샵 안내", ACCENT_BLUE, False),
        ("09:20", "10:10", "Guided 실습", "사전 주제로 PATH Step 1~4", ACCENT_PURPLE, False),
        ("10:10", "10:25", "휴식", "", CARD_BORDER, True),
        ("10:25", "11:30", "Guided 실습", "AI-DLC 코드 구현", ACCENT_PURPLE, False),
        ("11:30", "11:50", "자유 주제 준비", "브레인스토밍 + 설계", ACCENT_ORANGE, False),
        ("11:50", "12:50", "점심", "", CARD_BORDER, True),
        ("12:50", "14:20", "자유 실습", "본인 주제 PATH + 코드 구현", ACCENT_CYAN, False),
        ("14:20", "15:00", "결과 공유", "성과 발표 & 피드백", ACCENT_PINK, False),
    ]

    # Timeline card background
    tl_height = len(timeline) * Inches(0.38) + Inches(0.3)
    add_rrect(slide, MG, y, CW, tl_height, DARK_CARD, CARD_BORDER)
    ty = y + Inches(0.15)

    for i, (start, end, title, desc, color, is_break) in enumerate(timeline):
        # Dot
        dot_x = MG + Inches(0.2)
        add_circle(slide, dot_x, ty + Inches(0.08), Inches(0.1), color)

        # Vertical line
        if i < len(timeline) - 1:
            add_rect(slide, dot_x + Inches(0.04), ty + Inches(0.18),
                     Inches(0.02), Inches(0.25), CARD_BORDER)

        # Time
        time_str = f"{start} - {end}"
        txt(slide, MG + Inches(0.42), ty + Inches(0.02), Inches(1.3), Inches(0.2),
            time_str, 8, LIGHT_GRAY if is_break else color, not is_break)

        # Title
        txt(slide, MG + Inches(1.8), ty + Inches(0.02), Inches(1.8), Inches(0.2),
            title, 10, LIGHT_GRAY if is_break else WHITE, not is_break)

        # Desc
        if desc:
            txt(slide, MG + Inches(3.7), ty + Inches(0.02), Inches(2.5), Inches(0.2),
                desc, 8, LIGHT_GRAY)

        ty += Inches(0.38)

    y += tl_height + Inches(0.3)

    # ══════════════════════════════════════
    # WORKSHOP INFO
    # ══════════════════════════════════════
    add_rrect(slide, MG, y, CW, Inches(1.55), DARK_CARD, ACCENT_BLUE)
    iy = y + Inches(0.18)

    info_items = [
        ("일시", "2026년 O월 O일 (O) 09:00 - 15:00"),
        ("장소", "OOO (오프라인 / 온라인 병행)"),
        ("대상", "AI Agent 도입을 검토 중인 기업 담당자"),
        ("인원", "최대 20명 (선착순)"),
        ("준비물", "노트북 (AWS 계정 사전 안내)"),
    ]

    for label, value in info_items:
        txt(slide, MG + Inches(0.25), iy, Inches(0.6), Inches(0.2), label, 9, ACCENT_BLUE, True)
        txt(slide, MG + Inches(0.9), iy, Inches(4.8), Inches(0.2), value, 9, WHITE)
        iy += Inches(0.25)

    y += Inches(1.7)

    # CTA
    cta_w = Inches(3.2)
    cta_x = int((SW - cta_w) / 2)
    cta = add_rrect(slide, cta_x, y, cta_w, Inches(0.48), ACCENT_BLUE)
    badge_text(cta, "지금 신청하기 >", 14, WHITE)

    y += Inches(0.6)
    txt(slide, MG, y, CW, Inches(0.2), "문의: workshop@example.com", 8, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    y += Inches(0.3)
    add_rect(slide, 0, y, SW, Inches(0.015), CARD_BORDER)
    y += Inches(0.1)
    txt(slide, MG, y, CW, Inches(0.15),
        "P.A.T.H Agent Designer Workshop  |  Powered by AWS Bedrock + Strands Agents SDK",
        7, CARD_BORDER, False, PP_ALIGN.CENTER)

    # ─── Save ───
    out = "/home/ec2-user/project/path/docs/PATH-workshop-EDM.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    create_edm()
